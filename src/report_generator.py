"""
Report Generator Module
=======================
Automated Excel and PowerPoint report generation.

Author: Risk Analytics Team
Date: November 2025
"""

import pandas as pd
import numpy as np
import logging
from typing import Dict, List
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, LineChart, Reference

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class ReportGenerator:
    """
    Automated report generator for VaR analysis.
    
    Generates professional Excel workbooks and PowerPoint presentations
    with VaR results, backtesting analysis, and stress test outputs.
    
    Attributes:
        returns (pd.DataFrame): Historical returns
        weights (list): Portfolio weights
        portfolio_value (float): Total portfolio value
    """
    
    def __init__(self, returns: pd.DataFrame, weights: List[float], 
                 portfolio_value: float = 10_000_000):
        """
        Initialize report generator.
        
        Args:
            returns (pd.DataFrame): Historical returns
            weights (List[float]): Portfolio weights
            portfolio_value (float): Portfolio value
        """
        self.returns = returns
        self.weights = weights
        self.portfolio_value = portfolio_value
        self.report_date = datetime.now()
        
        logger.info("Report Generator initialized")
    
    def generate_excel_report(self, filepath: str, var_results: Dict, 
                            backtest_results: Dict, stress_results: Dict) -> None:
        """
        Generate comprehensive Excel report.
        
        Args:
            filepath (str): Output file path
            var_results (Dict): VaR calculation results
            backtest_results (Dict): Backtesting results
            stress_results (Dict): Stress test results
        """
        logger.info(f"Generating Excel report: {filepath}")
        
        # Create Excel writer
        with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
            # Sheet 1: Executive Summary
            self._create_summary_sheet(writer, var_results, backtest_results)
            
            # Sheet 2: VaR Calculations
            self._create_var_sheet(writer, var_results)
            
            # Sheet 3: Backtesting Results
            self._create_backtest_sheet(writer, backtest_results)
            
            # Sheet 4: Stress Testing
            self._create_stress_sheet(writer, stress_results)
            
            # Sheet 5: Portfolio Details
            self._create_portfolio_sheet(writer)
            
            # Sheet 6: Raw Data
            self._create_data_sheet(writer)
        
        # Apply formatting
        self._format_excel_report(filepath)
        
        logger.info(f"Excel report generated successfully: {filepath}")
    
    def _create_summary_sheet(self, writer, var_results: Dict, backtest_results: Dict) -> None:
        """Create executive summary sheet."""
        summary_data = {
            'Metric': [
                'Report Date',
                'Portfolio Value',
                'Number of Assets',
                'Observation Period',
                '',
                'VaR (95%, 1-day)',
                'VaR (99%, 1-day)',
                '',
                'Backtest Period',
                'Exceptions (99%)',
                'Exception Rate',
                'Traffic Light Status',
                '',
                'Model Assessment'
            ],
            'Value': [
                self.report_date.strftime('%Y-%m-%d'),
                f"${self.portfolio_value:,.0f}",
                len(self.weights),
                f"{len(self.returns)} days",
                '',
                f"${var_results.get('var_95', 0):,.2f}" if 'var_95' in var_results else 'N/A',
                f"${var_results.get('var_99', 0):,.2f}" if 'var_99' in var_results else 'N/A',
                '',
                f"{backtest_results['summary']['n_observations']} days" if backtest_results else 'N/A',
                backtest_results['summary']['n_exceptions'] if backtest_results else 'N/A',
                f"{backtest_results['summary']['exception_rate']*100:.2f}%" if backtest_results else 'N/A',
                backtest_results['summary']['traffic_light'] if backtest_results else 'N/A',
                '',
                'Model is performing adequately' if backtest_results and backtest_results['summary']['traffic_light'] == 'GREEN' else 'Review required'
            ]
        }
        
        summary_df = pd.DataFrame(summary_data)
        summary_df.to_excel(writer, sheet_name='Executive Summary', index=False)
    
    def _create_var_sheet(self, writer, var_results: Dict) -> None:
        """Create VaR calculations sheet."""
        # VaR estimates table
        var_data = []
        for conf in [0.95, 0.99]:
            key = f'var_{int(conf*100)}'
            if key in var_results:
                var_data.append({
                    'Confidence Level': f'{conf*100:.0f}%',
                    'VaR (Currency)': var_results[key],
                    'VaR (%)': (var_results[key] / self.portfolio_value) * 100,
                    'Time Horizon': '1 day'
                })
        
        if var_data:
            var_df = pd.DataFrame(var_data)
            var_df.to_excel(writer, sheet_name='VaR Calculations', index=False, startrow=0)
        
        # Portfolio statistics
        portfolio_returns = (self.returns * self.weights).sum(axis=1)
        stats_data = {
            'Statistic': [
                'Mean Daily Return',
                'Daily Volatility',
                'Annualized Return',
                'Annualized Volatility',
                'Sharpe Ratio',
                'Skewness',
                'Kurtosis',
                'Min Return',
                'Max Return',
                '5th Percentile',
                '95th Percentile'
            ],
            'Value': [
                f'{portfolio_returns.mean()*100:.4f}%',
                f'{portfolio_returns.std()*100:.4f}%',
                f'{portfolio_returns.mean()*252*100:.2f}%',
                f'{portfolio_returns.std()*np.sqrt(252)*100:.2f}%',
                f'{(portfolio_returns.mean()*252)/(portfolio_returns.std()*np.sqrt(252)):.2f}',
                f'{portfolio_returns.skew():.4f}',
                f'{portfolio_returns.kurtosis():.4f}',
                f'{portfolio_returns.min()*100:.4f}%',
                f'{portfolio_returns.max()*100:.4f}%',
                f'{portfolio_returns.quantile(0.05)*100:.4f}%',
                f'{portfolio_returns.quantile(0.95)*100:.4f}%'
            ]
        }
        
        stats_df = pd.DataFrame(stats_data)
        stats_df.to_excel(writer, sheet_name='VaR Calculations', index=False, startrow=len(var_data)+3)
    
    def _create_backtest_sheet(self, writer, backtest_results: Dict) -> None:
        """Create backtesting results sheet."""
        if not backtest_results:
            return
        
        # Summary statistics
        summary = backtest_results['summary']
        summary_data = {
            'Metric': [
                'Test Period',
                'Number of Observations',
                'Confidence Level',
                'Expected Exceptions',
                'Actual Exceptions',
                'Exception Rate',
                'Expected Rate',
                'Traffic Light Status',
                '',
                'Average VaR',
                'Maximum VaR',
                'Average Excess Loss',
                'Maximum Excess Loss'
            ],
            'Value': [
                f"{summary['start_date']} to {summary['end_date']}",
                summary['n_observations'],
                f"{summary['confidence']*100:.0f}%",
                f"{summary['expected_exceptions']:.1f}",
                summary['n_exceptions'],
                f"{summary['exception_rate']*100:.2f}%",
                f"{summary['expected_rate']*100:.2f}%",
                summary['traffic_light'],
                '',
                f"${summary['avg_var']:,.2f}",
                f"${summary['max_var']:,.2f}",
                f"${summary['avg_excess']:,.2f}",
                f"${summary['max_excess']:,.2f}"
            ]
        }
        
        summary_df = pd.DataFrame(summary_data)
        summary_df.to_excel(writer, sheet_name='Backtesting', index=False, startrow=0)
        
        # Exception details
        if summary['exception_dates']:
            exception_data = pd.DataFrame({
                'Date': summary['exception_dates'],
                'Loss': summary['exception_losses'],
                'Excess': summary['excess_losses']
            })
            exception_data.to_excel(writer, sheet_name='Backtesting', index=False, startrow=len(summary_data)+3)
    
    def _create_stress_sheet(self, writer, stress_results: Dict) -> None:
        """Create stress testing sheet."""
        if not stress_results or 'standard_scenarios' not in stress_results:
            return
        
        # Standard scenarios
        standard_data = []
        for scenario_name, result in stress_results['standard_scenarios'].items():
            if 'portfolio_loss' in result:
                standard_data.append({
                    'Scenario': scenario_name,
                    'Portfolio Loss': result['portfolio_loss'],
                    'Loss (%)': result['portfolio_loss_pct']
                })
        
        if standard_data:
            standard_df = pd.DataFrame(standard_data)
            standard_df.to_excel(writer, sheet_name='Stress Testing', index=False, startrow=0)
        
        # Individual shocks
        if 'individual_shocks' in stress_results:
            individual_data = []
            for scenario_name, result in stress_results['individual_shocks'].items():
                if 'portfolio_loss' in result:
                    individual_data.append({
                        'Scenario': scenario_name,
                        'Portfolio Loss': result['portfolio_loss'],
                        'Loss (%)': result['portfolio_loss_pct']
                    })
            
            if individual_data:
                individual_df = pd.DataFrame(individual_data)
                individual_df.to_excel(writer, sheet_name='Stress Testing', 
                                      index=False, startrow=len(standard_data)+3)
    
    def _create_portfolio_sheet(self, writer) -> None:
        """Create portfolio details sheet."""
        portfolio_data = {
            'Asset': self.returns.columns.tolist(),
            'Weight': self.weights,
            'Weight (%)': [w*100 for w in self.weights],
            'Value': [w*self.portfolio_value for w in self.weights]
        }
        
        portfolio_df = pd.DataFrame(portfolio_data)
        portfolio_df.to_excel(writer, sheet_name='Portfolio Details', index=False)
    
    def _create_data_sheet(self, writer) -> None:
        """Create raw data sheet."""
        self.returns.to_excel(writer, sheet_name='Raw Returns')
    
    def _format_excel_report(self, filepath: str) -> None:
        """Apply formatting to Excel report."""
        wb = openpyxl.load_workbook(filepath)
        
        # Define styles
        header_fill = PatternFill(start_color='366092', end_color='366092', fill_type='solid')
        header_font = Font(color='FFFFFF', bold=True)
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Format headers
            for cell in ws[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal='center', vertical='center')
            
            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column = [cell for cell in column]
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column[0].column_letter].width = adjusted_width
        
        wb.save(filepath)
        logger.info("Excel formatting applied")
    
    def generate_ppt_report(self, filepath: str, var_results: Dict, 
                           backtest_results: Dict, stress_results: Dict,
                           chart_dir: str = 'reports/charts') -> None:
        """
        Generate PowerPoint presentation.
        
        Args:
            filepath (str): Output file path
            var_results (Dict): VaR results
            backtest_results (Dict): Backtesting results
            stress_results (Dict): Stress test results
            chart_dir (str): Directory containing chart images
        """
        logger.info(f"Generating PowerPoint report: {filepath}")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title
        self._add_title_slide(prs)
        
        # Slide 2: Executive Summary
        self._add_summary_slide(prs, var_results, backtest_results)
        
        # Slide 3: VaR Methodology
        self._add_methodology_slide(prs)
        
        # Slide 4: VaR Results
        self._add_var_results_slide(prs, var_results)
        
        # Slide 5: Backtesting Analysis
        self._add_backtest_slide(prs, backtest_results)
        
        # Slide 6: Stress Testing
        self._add_stress_slide(prs, stress_results)
        
        # Slide 7: Recommendations
        self._add_recommendations_slide(prs, backtest_results)
        
        prs.save(filepath)
        logger.info(f"PowerPoint report generated successfully: {filepath}")
    
    def _add_title_slide(self, prs) -> None:
        """Add title slide."""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Market Risk VaR Analysis"
        subtitle.text = f"Value-at-Risk Model Development & Monitoring\n{self.report_date.strftime('%B %d, %Y')}"
        
        # Format title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(54, 96, 146)
    
    def _add_summary_slide(self, prs, var_results: Dict, backtest_results: Dict) -> None:
        """Add executive summary slide."""
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Add content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        content = f"""Portfolio Overview:
• Portfolio Value: ${self.portfolio_value:,.0f}
• Number of Assets: {len(self.weights)}
• Analysis Period: {len(self.returns)} trading days

Risk Metrics:
• 95% VaR (1-day): ${var_results.get('var_95', 0):,.2f}
• 99% VaR (1-day): ${var_results.get('var_99', 0):,.2f}

Model Validation:
• Backtesting Period: {backtest_results['summary']['n_observations'] if backtest_results else 'N/A'} days
• Exception Rate: {backtest_results['summary']['exception_rate']*100:.2f}% if backtest_results else 'N/A'}
• Traffic Light: {backtest_results['summary']['traffic_light'] if backtest_results else 'N/A'}
• Model Status: {'ADEQUATE' if backtest_results and backtest_results['summary']['traffic_light'] == 'GREEN' else 'REVIEW REQUIRED'}"""
        
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(14)
    
    def _add_methodology_slide(self, prs) -> None:
        """Add methodology slide."""
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title = slide.shapes.title
        title.text = "VaR Methodology"
        
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        content = """Historical Simulation VaR:
• Uses actual historical return distribution
• Non-parametric approach (no distributional assumptions)
• Calculates VaR as percentile of historical losses
• 95% VaR: Loss exceeded 5% of the time
• 99% VaR: Loss exceeded 1% of the time

Backtesting Framework:
• Compares predicted VaR vs actual daily losses
• Counts exceptions (losses exceeding VaR)
• Basel traffic-light classification:
  - Green Zone: 0-4 exceptions (99% VaR, 250 days)
  - Yellow Zone: 5-9 exceptions - monitoring required
  - Red Zone: 10+ exceptions - model inadequate"""
        
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(13)
    
    def _add_var_results_slide(self, prs, var_results: Dict) -> None:
        """Add VaR results slide."""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.75)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_box.text = "VaR Results"
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        title_box.text_frame.paragraphs[0].font.bold = True
        
        # VaR values
        left = Inches(1)
        top = Inches(2)
        
        content_box = slide.shapes.add_textbox(left, top, Inches(8), Inches(4))
        tf = content_box.text_frame
        
        content = f"""1-Day VaR Estimates:

95% Confidence Level:
• Maximum potential loss: ${var_results.get('var_95', 0):,.2f}
• Probability: Losses will not exceed this 95% of days

99% Confidence Level:
• Maximum potential loss: ${var_results.get('var_99', 0):,.2f}
• Probability: Losses will not exceed this 99% of days

Interpretation: We are 99% confident that portfolio losses will not exceed 
${var_results.get('var_99', 0):,.2f} over the next trading day."""
        
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(16)
    
    def _add_backtest_slide(self, prs, backtest_results: Dict) -> None:
        """Add backtesting slide."""
        if not backtest_results:
            return
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_box.text = "Backtesting Analysis"
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        title_box.text_frame.paragraphs[0].font.bold = True
        
        # Results
        summary = backtest_results['summary']
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = content_box.text_frame
        
        status_emoji = "✓" if summary['traffic_light'] == 'GREEN' else "⚠" if summary['traffic_light'] == 'YELLOW' else "✗"
        
        content = f"""Validation Results (99% VaR):

Test Period: {summary['n_observations']} trading days
Exceptions: {summary['n_exceptions']} ({summary['exception_rate']*100:.2f}%)
Expected: {summary['expected_exceptions']:.1f} exceptions

Traffic Light Classification: {summary['traffic_light']} {status_emoji}

Model Assessment:
{self._get_backtest_interpretation(summary['traffic_light'])}"""
        
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(16)
    
    def _add_stress_slide(self, prs, stress_results: Dict) -> None:
        """Add stress testing slide."""
        if not stress_results:
            return
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_box.text = "Stress Testing Results"
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        title_box.text_frame.paragraphs[0].font.bold = True
        
        # Get worst scenarios
        if 'standard_scenarios' in stress_results:
            scenarios = []
            for name, result in stress_results['standard_scenarios'].items():
                if 'portfolio_loss' in result:
                    scenarios.append((name, result['portfolio_loss']))
            
            scenarios.sort(key=lambda x: x[1], reverse=True)
            top_scenarios = scenarios[:4]
            
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = content_box.text_frame
            
            content = "Scenario Losses:\n\n"
            for name, loss in top_scenarios:
                content += f"• {name}: ${loss:,.2f}\n"
            
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(16)
    
    def _add_recommendations_slide(self, prs, backtest_results: Dict) -> None:
        """Add recommendations slide."""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_box.text = "Recommendations"
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        title_box.text_frame.paragraphs[0].font.bold = True
        
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = content_box.text_frame
        
        if backtest_results and backtest_results['summary']['traffic_light'] == 'GREEN':
            recommendations = """✓ Model Performance: SATISFACTORY

Current Actions:
• Continue daily VaR monitoring
• Maintain current risk limits
• Update model quarterly with new data

Best Practices:
• Monitor for market regime changes
• Conduct monthly backtesting
• Review stress scenarios quarterly
• Document all model changes"""
        else:
            recommendations = """⚠ Model Performance: REQUIRES ATTENTION

Immediate Actions:
• Increase backtesting frequency to daily
• Review recent exceptions for patterns
• Consider recalibrating model parameters
• Implement additional risk controls

Investigation Required:
• Analyze market conditions during exceptions
• Assess if portfolio composition has changed
• Evaluate if volatility has shifted
• Consider alternative VaR methodologies"""
        
        p = tf.paragraphs[0]
        p.text = recommendations
        p.font.size = Pt(14)
    
    def _get_backtest_interpretation(self, traffic_light: str) -> str:
        """Get interpretation text for traffic light status."""
        interpretations = {
            'GREEN': "Model is accurately predicting risk. Performance is within acceptable limits.",
            'YELLOW': "Model requires monitoring. Exception rate is elevated but still within tolerance.",
            'RED': "Model is inadequate. Immediate review and recalibration required."
        }
        return interpretations.get(traffic_light, "Status unknown")


# Example usage
if __name__ == "__main__":
    logger.info("Report Generator module ready")
